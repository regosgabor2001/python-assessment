import unittest
import json
from unittest.mock import patch
from pptx.presentation import Presentation as PresentationClass
from report_generator import readFile, createPresentation

class TestPresentationCreation(unittest.TestCase):
    def test_createPresentation(self):
        data = readFile()

        presentation = createPresentation(data)

        #assert that the presentation is an instance of the PresentationClass
        self.assertIsInstance(presentation, PresentationClass)

        #assert the number of slides in the presentation
        self.assertEqual(len(presentation.slides), 5)


        #check the title and subtitle of the slides
        title_slide = presentation.slides[0]
        self.assertEqual(title_slide.shapes.title.text, 'The Title Text')
        self.assertEqual(title_slide.placeholders[1].text, 'The Sub-Title Text')

        text_slide = presentation.slides[1]
        self.assertEqual(text_slide.shapes.title.text, 'The Title Text')
        self.assertEqual(text_slide.shapes[1].text_frame.text, 'The Long Text')

        text_slide = presentation.slides[2]
        self.assertEqual(text_slide.shapes.title.text, 'The Title Text')
        self.assertEqual(text_slide.shapes[1].text_frame.text, '\nThe Level 1 Text\nThe Level 2 Text\nThe Level 2 Text\nThe Level 1 Text')


if __name__ == '__main__':
    unittest.main()