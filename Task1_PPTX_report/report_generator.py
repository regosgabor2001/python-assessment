import json
import csv
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import XyChartData
import math
import logging

#sets up the logging configuration, specifying the filename, log level (INFO), and format of log messages
logging.basicConfig(filename='logger.log', level=logging.INFO, format='%(asctime)s:%(levelname)s:%(message)s')

#reads the content of the sample.json file and loads it as a Python dictionary using the json.load() function. The function then returns the loaded data, the function then returns the loaded data
def readFile():
    with open("Task1_PPTX_report/sample.json", "r") as f:
        data = json.load(f)
    return data

def createPresentation(data):
    presentation = Presentation()

    for item in data['presentation']:
        if item['type'] == 'title':
            #create a slide with a title and subtitle
            slide = presentation.slides.add_slide(presentation.slide_layouts[0])

            #access the title and subtitle shapes on the slide
            title = slide.shapes.title
            subtitle = slide.placeholders[1]

            #set the text for the title and subtitle
            title.text = item['title']
            subtitle.text = item['content']

        elif item['type'] == 'text':
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            title = slide.shapes.title
            content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(1), Inches(1)).text_frame

            title.text = item['title']
            content.text = item['content']

        elif item['type'] == 'list':
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            title = slide.shapes.title

            title.text = item['title']

            for thing in item['content']:
                if thing['level'] == 1:
                    #add a paragraph with level 1 indentation to the content placeholder
                    content = slide.shapes.placeholders[1].text_frame.add_paragraph()
                    content.text = thing['text']
                    content.level = 1
                elif thing['level'] == 2:
                    #add a paragraph with level 2 indentation to the content placeholder
                    content = slide.shapes.placeholders[1].text_frame.add_paragraph()
                    content.text = thing['text']
                    content.level = 2

            
        elif item['type'] == 'picture':
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            title = slide.shapes.title

            #get the path or URL of the picture from the current item
            pic = item['content']

            title.text = item['title']
            
            #add the picture to the slide at the specified position
            slide.shapes.add_picture(pic,Inches(3),Inches(2))

        elif item['type'] == 'plot':
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            title = slide.shapes.title

            title.text = item['title']

            xValues = []
            yValues = []

            #read the data from the specified CSV file
            with open(item['content'], 'r') as f:
                csv_reader = csv.reader(f, delimiter=';')
            
                #extract x and y values from each row in the CSV file
                for row in csv_reader:
                    xValues.append(row[0])
                    yValues.append(row[1])

            #create the chart data object
            chart_data = XyChartData()
            cd = chart_data.add_series('',number_format=None)

            #add data points to the chart data series
            for x, y in list(zip(xValues, yValues)):
                cd.add_data_point(x, y, number_format=None)


            #add the chart to the slide at the specified position and size
            x, y, cx, cy = Inches(1), Inches(2), Inches(7), Inches(5)
            chart = slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS, x, y, cx, cy, chart_data).chart

            #set the axis titles based on the configuration values from the item
            chart.category_axis.axis_title.text_frame.text= item['configuration']['x-label']
            chart.value_axis.axis_title.text_frame.text= item['configuration']['y-label']

            #customize chart settings
            chart.has_legend = False
            for series in chart.series:
                series.format.line.visible = False

            #set the minimum and maximum scale values for the axes
            chart.value_axis.minimum_scale = math.floor(float(yValues[0]))
            chart.category_axis.minimum_scale = math.floor(float(xValues[0]))

            chart.value_axis.maximum_scale = math.ceil(float(yValues[len(yValues)-1]))
            chart.category_axis.maximum_scale = math.ceil(float(xValues[len(xValues)-1]))

            


    return presentation


data = readFile()
presentation = createPresentation(data)

#takes the created presentation as input and saves it as "result.pptx"
presentation.save("result.pptx")

#try to execute the following code block, and catch any exceptions that occure
try:
    #log information messages
    logging.info("Reading data from file...")
    data = readFile()
    logging.info("Creating presentation...")
    presentation = createPresentation(data)
    logging.info("Saving presentation...")
    presentation.save("result.pptx")
    logging.info("Presentation saved successfully.")
except Exception as e:
    #log an error message with the details of the exception
    logging.error(f"An error occurred: {str(e)}")